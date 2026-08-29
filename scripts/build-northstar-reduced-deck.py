#!/usr/bin/env python3
"""NorthStar reduced deck — 2 slides: path+confirm, this week+lock."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = "interview-prep/ramp/northstar-business-case/NorthStar-Reduced-2-Slides.pptx"

NAVY = RGBColor(0x1B, 0x28, 0x38)
SLATE = RGBColor(0x4A, 0x55, 0x68)
MIST = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xC9, 0x8A, 0x2E)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
LIGHT_TEAL = RGBColor(0xE6, 0xF4, 0xF4)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)

W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
blank = prs.slide_layouts[6]


def send_back(slide, shape):
    spTree = slide.shapes._spTree
    el = shape._element
    spTree.remove(el)
    spTree.insert(2, el)


def bg(slide, color=MIST):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    send_back(slide, s)


def tb(slide, left, top, width, height, lines, size=14, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold if i == 0 else False
        # allow per-call bold for all lines when single emphasis block
        if bold and len(lines) > 1 and i > 0:
            p.font.bold = True
        p.font.color.rgb = color
        p.alignment = align
        p.space_after = Pt(2)
    return box


def tb_all_bold(slide, left, top, width, height, lines, size=14, color=NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = align
    return box


# ── SLIDE 1: Path + confirm ──
s1 = prs.slides.add_slide(blank)
bg(s1)

# thin top bar
bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = NAVY
bar.line.fill.background()

tb(s1, Inches(0.6), Inches(0.28), Inches(10), Inches(0.35),
   "NorthStar Aviation  ·  Aircraft parts path", size=12, color=SLATE)
tb(s1, Inches(0.6), Inches(0.55), Inches(12), Inches(0.5),
   "Recommended path — pressure-test with me", size=28, bold=True, color=NAVY)

steps = [
    ("1", "Request", ["Qty × price", "Need-by · vendor", "Part #"]),
    ("2", "Approve", ["Your thresholds", "Inventory? (TBD)"]),
    ("3", "PO", ["On approval", "→ NetSuite"]),
    ("4", "Receive", ["Where you", "receive today"]),
    ("5", "Match & pay", ["3-way match", "Bill Pay closes", "the loop"]),
]
x0, sw, gap, top = Inches(0.5), Inches(2.4), Inches(0.1), Inches(1.25)
for i, (num, title, lines) in enumerate(steps):
    left = x0 + i * (sw + gap)
    box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, sw, Inches(2.55))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_TEAL
    box.line.color.rgb = BORDER
    circ = s1.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.12), top + Inches(0.15), Inches(0.36), Inches(0.36))
    circ.fill.solid()
    circ.fill.fore_color.rgb = TEAL if i < 4 else ACCENT
    circ.line.fill.background()
    tb(s1, left + Inches(0.12), top + Inches(0.17), Inches(0.36), Inches(0.32),
       num, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s1, left + Inches(0.12), top + Inches(0.6), sw - Inches(0.24), Inches(0.4),
       title, size=15, bold=True, color=NAVY)
    tb(s1, left + Inches(0.12), top + Inches(1.05), sw - Inches(0.24), Inches(1.35),
       lines, size=12, color=SLATE)

tb(s1, Inches(0.6), Inches(3.95), Inches(12), Inches(0.3),
   "Parts = three-way (order + receipt + invoice). NetSuite stays the books. Services later.", size=12, bold=True, color=TEAL)

# Confirm panel
panel = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.35), Inches(12.35), Inches(2.7))
panel.fill.solid()
panel.fill.fore_color.rgb = WHITE
panel.line.color.rgb = ACCENT
panel.line.width = Pt(2.5)

tb(s1, Inches(0.8), Inches(4.5), Inches(8), Inches(0.35),
   "Confirm with me", size=16, bold=True, color=ACCENT)

prompts = [
    "Inventory in the path — visibility, accountability, or both?  ________________",
    "Approval thresholds — who / at what $?  ________________",
    "Receiving — system + who? (Requester ≠ sole receiver)  ________________",
    "Pilot group + working session time  ________________",
]
for j, line in enumerate(prompts):
    tb(s1, Inches(0.8), Inches(5.0 + j * 0.45), Inches(11.8), Inches(0.4), line, size=14, color=NAVY)

# footer
tb(s1, Inches(0.6), Inches(7.15), Inches(6), Inches(0.25), "1 / 2  ·  Point · ask · write", size=10, color=SLATE)

# ── SLIDE 2: This week + lock ──
s2 = prs.slides.add_slide(blank)
bg(s2)
bar2 = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
bar2.fill.solid()
bar2.fill.fore_color.rgb = NAVY
bar2.line.fill.background()

tb(s2, Inches(0.6), Inches(0.35), Inches(12), Inches(0.5),
   "This week — first path live", size=28, bold=True, color=NAVY)

hero = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.05), Inches(12.35), Inches(1.15))
hero.fill.solid()
hero.fill.fore_color.rgb = NAVY
hero.line.fill.background()
tb(s2, Inches(0.85), Inches(1.25), Inches(11.8), Inches(0.35),
   "Aircraft Parts Program  ·  small pilot  ·  one real request end-to-end", size=18, bold=True, color=WHITE)
tb(s2, Inches(0.85), Inches(1.7), Inches(11.8), Inches(0.35),
   "Prove the path. Then expand. Don’t design the whole company day one.", size=13, color=RGBColor(0xCB, 0xD5, 0xE1))

actions = [
    ("1", "Build & publish\nthe Parts Program"),
    ("2", "Run the first\nrequest together"),
    ("3", "Review what broke\n— then scale"),
]
for k, (num, label) in enumerate(actions):
    left = Inches(0.5) + k * Inches(4.15)
    box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.5), Inches(3.95), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = BORDER
    tb(s2, left + Inches(0.2), Inches(2.65), Inches(0.5), Inches(0.5), num, size=26, bold=True, color=TEAL)
    lines = label.split("\n")
    tb(s2, left + Inches(0.75), Inches(2.75), Inches(3.0), Inches(1.1), lines, size=14, bold=True, color=NAVY)

lock = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.3), Inches(12.35), Inches(2.4))
lock.fill.solid()
lock.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
lock.line.color.rgb = ACCENT
lock.line.width = Pt(2.5)

tb(s2, Inches(0.8), Inches(4.45), Inches(8), Inches(0.35),
   "Lock before we leave", size=16, bold=True, color=ACCENT)

for m, item in enumerate([
    "☐  Inventory decision for the pilot",
    "☐  Approval thresholds (or owner who has them)",
    "☐  Receiving owner + system",
    "☐  Pilot group + working session on the calendar",
]):
    tb(s2, Inches(0.8), Inches(4.95 + m * 0.4), Inches(11.5), Inches(0.38), item, size=14, color=NAVY)

tb(s2, Inches(0.6), Inches(7.15), Inches(10), Inches(0.25),
   "2 / 2  ·  First milestone: program live + first request through Ramp", size=10, color=SLATE)

prs.save(OUT)
print("Saved", OUT)
