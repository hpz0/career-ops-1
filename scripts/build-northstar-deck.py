#!/usr/bin/env python3
"""Build NorthStar Aviation design-session deck (5 slides)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = "interview-prep/ramp/NorthStar-Aviation-Procurement-Design-Session.pptx"

NAVY = RGBColor(0x1B, 0x28, 0x38)
SLATE = RGBColor(0x4A, 0x55, 0x68)
MIST = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xC9, 0x8A, 0x2E)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
LIGHT_TEAL = RGBColor(0xE6, 0xF4, 0xF4)
CARD_BG = RGBColor(0xF8, 0xFA, 0xFC)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)

W, H = Inches(13.333), Inches(7.5)

NOTES = {
    1: "Open: expense live → design purchasing. Two pains: speed + audit. 15 min to align on parts path + lock first step.",
    2: "Do not read cards. Ask: Does one parts path match your priority? Regional growth = say once, not from slide.",
    3: "Walk the 5 steps in 60 sec. THEN Inventory question — visibility vs accountability vs both. Write her answer on notepad. Three-way = define once if she asks.",
    4: "90 seconds max. Path first. AI does not approve — optional later. Do not read vendor research line unless she asks.",
    5: "Close: calendar hold. Read lock checklist aloud and assign owners. 60-day win = one trusted path.",
}

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
blank = prs.slide_layouts[6]


def send_to_back(slide, shape):
    spTree = slide.shapes._spTree
    el = shape._element
    spTree.remove(el)
    spTree.insert(2, el)


def add_bg(slide, color=MIST):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    send_to_back(slide, bg)


def add_footer(slide, section):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, H - Inches(0.35), W, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xE8, 0xED, 0xF2)
    bar.line.fill.background()
    tf = slide.shapes.add_textbox(Inches(0.5), H - Inches(0.32), Inches(6), Inches(0.28)).text_frame
    p = tf.paragraphs[0]
    p.text = section
    p.font.size = Pt(9)
    p.font.color.rgb = SLATE
    tf2 = slide.shapes.add_textbox(Inches(11.2), H - Inches(0.32), Inches(1.8), Inches(0.28)).text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "NorthStar Aviation"
    p2.font.size = Pt(9)
    p2.font.color.rgb = SLATE
    p2.alignment = PP_ALIGN.RIGHT


def write_text(tf, lines, size=14, bold=False, color=NAVY, align=PP_ALIGN.LEFT, spacing=1.15):
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold and i == 0
        p.font.color.rgb = color
        p.alignment = align
        p.space_after = Pt(4)
        p.line_spacing = spacing


def textbox(slide, left, top, width, height, lines, size=14, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if isinstance(lines, str):
        lines = [lines]
    write_text(tf, lines, size=size, bold=bold, color=color, align=align)
    return tb


def card(slide, left, top, width, height, title, body_lines):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.07), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    textbox(slide, left + Inches(0.22), top + Inches(0.18), width - Inches(0.35), Inches(0.4), title, size=17, bold=True)
    textbox(slide, left + Inches(0.22), top + Inches(0.62), width - Inches(0.35), height - Inches(0.7), body_lines, size=13, color=SLATE)


def add_notes(slide, note):
    notes = slide.notes_slide
    notes.notes_text_frame.text = note


# ── Slide 1 ──
s1 = prs.slides.add_slide(blank)
add_bg(s1, NAVY)
textbox(s1, Inches(0.85), Inches(2.05), Inches(11.5), Inches(1.0), "Procurement Design & Activation", size=38, bold=True, color=WHITE)
textbox(s1, Inches(0.85), Inches(3.15), Inches(11), Inches(0.55),
        "NorthStar Aviation  ·  Design session with Monica (Controller)", size=19, color=RGBColor(0xCB, 0xD5, 0xE1))
textbox(s1, Inches(0.85), Inches(3.95), Inches(11), Inches(0.45), "Align on design. Lock the first step.", size=18, color=ACCENT)
textbox(s1, Inches(0.85), Inches(5.45), Inches(11), Inches(0.35), "Harrison Pizzi  ·  August 2026", size=13, color=RGBColor(0x94, 0xA3, 0xB8))
add_notes(s1, NOTES[1])

# ── Slide 2 ──
s2 = prs.slides.add_slide(blank)
add_bg(s2, MIST)
add_footer(s2, "What we're solving")
textbox(s2, Inches(0.7), Inches(0.4), Inches(8), Inches(0.55), "What we're solving", size=30, bold=True)

card(s2, Inches(0.7), Inches(1.2), Inches(5.85), Inches(2.0), "Approval speed",
     ["Approvals bottleneck.", "Parts miss need-by dates."])
card(s2, Inches(6.78), Inches(1.2), Inches(5.85), Inches(2.0), "Audit readiness",
     ["No single trail request → payment.", "Email, Slack, spreadsheets."])

banner = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(3.55), Inches(11.93), Inches(1.65))
banner.fill.solid()
banner.fill.fore_color.rgb = TEAL
banner.line.fill.background()
textbox(s2, Inches(1.0), Inches(3.75), Inches(11), Inches(0.35), "START HERE", size=11, bold=True, color=RGBColor(0xB2, 0xDF, 0xDB))
textbox(s2, Inches(1.0), Inches(4.15), Inches(11.2), Inches(0.95),
        ["One controlled aircraft parts path — request through payment.",
         "Prove it with a pilot, then expand."], size=19, bold=True, color=WHITE)
add_notes(s2, NOTES[2])

# ── Slide 3 ──
s3 = prs.slides.add_slide(blank)
add_bg(s3, MIST)
add_footer(s3, "Recommended path")
textbox(s3, Inches(0.7), Inches(0.4), Inches(11), Inches(0.55), "Recommended path · confirm together", size=30, bold=True)

steps = [
    ("Request", ["Parts Program", "Qty × price · need-by", "Vendor · part #"]),
    ("Approve", ["Your thresholds", "Inventory? (TBD)"]),
    ("PO", ["On approval", "→ NetSuite"]),
    ("Receive", ["Where you", "receive today"]),
    ("Match & pay", ["Order + receipt", "+ invoice → pay"]),
]
x0, step_w, gap, top = Inches(0.55), Inches(2.38), Inches(0.1), Inches(1.15)
for i, (title, lines) in enumerate(steps):
    left = x0 + i * (step_w + gap)
    h = Inches(2.45)
    box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, step_w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_TEAL
    box.line.color.rgb = BORDER
    circ = s3.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.12), top + Inches(0.12), Inches(0.38), Inches(0.38))
    circ.fill.solid()
    circ.fill.fore_color.rgb = TEAL if i < 4 else ACCENT
    circ.line.fill.background()
    textbox(s3, left + Inches(0.12), top + Inches(0.14), Inches(0.38), Inches(0.32), str(i + 1), size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(s3, left + Inches(0.12), top + Inches(0.58), step_w - Inches(0.24), Inches(0.4), title, size=14, bold=True)
    textbox(s3, left + Inches(0.12), top + Inches(0.98), step_w - Inches(0.24), Inches(1.35), lines, size=11, color=SLATE)
    if i < 4:
        arr = s3.shapes.add_shape(MSO_SHAPE.CHEVRON, left + step_w + Inches(0.01), top + Inches(1.05), Inches(0.08), Inches(0.35))
        arr.fill.solid()
        arr.fill.fore_color.rgb = RGBColor(0x94, 0xA3, 0xB8)
        arr.line.fill.background()

textbox(s3, Inches(0.7), Inches(3.72), Inches(12), Inches(0.3),
        "Three-way = order + receipt + invoice before pay", size=12, bold=True, color=TEAL)

panel = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(4.15), Inches(11.93), Inches(2.45))
panel.fill.solid()
panel.fill.fore_color.rgb = WHITE
panel.line.color.rgb = ACCENT
panel.line.width = Pt(2)
textbox(s3, Inches(1.0), Inches(4.28), Inches(4), Inches(0.35), "Confirm with me", size=16, bold=True, color=ACCENT)
prompts = [
    "Inventory in the path — visibility, accountability, or both?  _______________",
    "Approval thresholds — who / at what $  _______________",
    "Receiving — system + owner  _______________",
    "Pilot group + working session (45–60 min)  _______________",
]
for j, line in enumerate(prompts):
    textbox(s3, Inches(1.0), Inches(4.78 + j * 0.45), Inches(11.3), Inches(0.4), line, size=13)
add_notes(s3, NOTES[3])

# ── Slide 4 ──
s4 = prs.slides.add_slide(blank)
add_bg(s4, MIST)
add_footer(s4, "After the path is live")
textbox(s4, Inches(0.7), Inches(0.4), Inches(10), Inches(0.55), "After the path is live", size=30, bold=True)
textbox(s4, Inches(0.7), Inches(0.98), Inches(10), Inches(0.35), "Extensions — not week-one blockers", size=15, color=SLATE)
card(s4, Inches(0.7), Inches(1.55), Inches(5.85), Inches(3.35), "Bill Pay",
     ["Matched invoices paid in Ramp.", "", "Due dates leave the spreadsheet.", "", "Full request-to-payment trail."])
card(s4, Inches(6.78), Inches(1.55), Inches(5.85), Inches(3.35), "AI summary (optional)",
     ["Short summary for human approvers.", "", "Does not approve or reject.", "", "Useful as volume grows — path first."])
textbox(s4, Inches(0.7), Inches(5.15), Inches(11.5), Inches(0.45),
        "Vendor research agents can help on new vendors later — not required to prove the parts path.", size=11, color=SLATE)
add_notes(s4, NOTES[4])

# ── Slide 5 ──
s5 = prs.slides.add_slide(blank)
add_bg(s5, MIST)
add_footer(s5, "This week")
textbox(s5, Inches(0.7), Inches(0.4), Inches(8), Inches(0.55), "This week", size=30, bold=True)
hero = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.1), Inches(11.93), Inches(1.05))
hero.fill.solid()
hero.fill.fore_color.rgb = NAVY
hero.line.fill.background()
textbox(s5, Inches(1.0), Inches(1.38), Inches(11.4), Inches(0.55),
        "Stand up the Aircraft Parts Program  ·  pilot group  ·  one real request end-to-end", size=18, bold=True, color=WHITE)

actions = [
    ("1", "Build & publish\n the Parts Program"),
    ("2", "Run the first\n request together"),
    ("3", "Review what broke\n — tune before scaling"),
]
for k, (num, label) in enumerate(actions):
    left = Inches(0.7) + k * Inches(4.05)
    box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.45), Inches(3.85), Inches(1.45))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = BORDER
    textbox(s5, left + Inches(0.18), Inches(2.58), Inches(0.55), Inches(0.55), num, size=26, bold=True, color=TEAL)
    textbox(s5, left + Inches(0.72), Inches(2.68), Inches(2.95), Inches(1.1), label.split("\n"), size=13, bold=True)

lock = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(4.15), Inches(11.93), Inches(2.15))
lock.fill.solid()
lock.fill.fore_color.rgb = CARD_BG
lock.line.color.rgb = ACCENT
lock.line.width = Pt(2)
textbox(s5, Inches(1.0), Inches(4.28), Inches(5), Inches(0.35), "Lock before we leave", size=16, bold=True, color=ACCENT)
for m, item in enumerate([
    "☐  Inventory decision for the pilot",
    "☐  Approval thresholds (or owner who has them)",
    "☐  Receiving owner + system",
    "☐  Pilot group + session on calendar",
]):
    textbox(s5, Inches(1.0), Inches(4.78 + m * 0.42), Inches(11), Inches(0.38), item, size=14)
textbox(s5, Inches(0.7), Inches(6.45), Inches(11.5), Inches(0.3),
        "60-day win = one trusted parts path — not every edge case on day one.", size=11, color=SLATE)
add_notes(s5, NOTES[5])

prs.save(OUT)
print("Saved:", OUT)
